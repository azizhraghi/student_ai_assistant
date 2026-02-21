"""
PowerPoint Parser Tool — extracts text from .pptx files using python-pptx.
"""

from pptx import Presentation
import io


def parse_pptx(file_bytes: bytes) -> str:
    """
    Extract all text from a PowerPoint file given its raw bytes.
    Returns a structured string with slide content.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        if slide_content:
            slides_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_content))

    return "\n\n".join(slides_text) if slides_text else "No text found in PowerPoint."
